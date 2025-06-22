import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

# Create test directory if not exists
os.makedirs("test-materials", exist_ok=True)

# 1. Create DOCX file
doc = Document()
doc.add_heading('데이터베이스 설계 원칙', 0)

doc.add_heading('1. 정규화(Normalization)', level=1)
doc.add_paragraph('데이터베이스 정규화는 데이터 중복을 최소화하고 데이터 무결성을 향상시키는 과정입니다.')
doc.add_paragraph('• 제1정규형(1NF): 모든 속성값은 원자값이어야 함')
doc.add_paragraph('• 제2정규형(2NF): 부분 함수 종속 제거')
doc.add_paragraph('• 제3정규형(3NF): 이행 함수 종속 제거')
doc.add_paragraph('• BCNF: 모든 결정자가 후보키')

doc.add_heading('2. 인덱스(Index)', level=1)
doc.add_paragraph('인덱스는 데이터베이스 검색 속도를 향상시키는 자료구조입니다.')
doc.add_paragraph('• B-Tree 인덱스: 가장 일반적인 인덱스 구조')
doc.add_paragraph('• Hash 인덱스: 동등 비교에 최적화')
doc.add_paragraph('• Bitmap 인덱스: 카디널리티가 낮은 컬럼에 적합')

doc.add_heading('3. 트랜잭션(Transaction)', level=1)
doc.add_paragraph('트랜잭션은 데이터베이스의 상태를 변화시키는 작업의 단위입니다.')
doc.add_paragraph('ACID 속성:')
doc.add_paragraph('• Atomicity(원자성): 모든 작업이 완전히 수행되거나 전혀 수행되지 않음')
doc.add_paragraph('• Consistency(일관성): 트랜잭션 전후로 데이터베이스의 일관성 유지')
doc.add_paragraph('• Isolation(격리성): 동시에 실행되는 트랜잭션들이 서로 영향을 주지 않음')
doc.add_paragraph('• Durability(지속성): 완료된 트랜잭션의 결과는 영구적으로 반영')

doc.add_heading('4. SQL 최적화', level=1)
doc.add_paragraph('• 적절한 인덱스 사용')
doc.add_paragraph('• 불필요한 조인 제거')
doc.add_paragraph('• 서브쿼리 대신 조인 사용')
doc.add_paragraph('• EXPLAIN 명령어로 실행 계획 확인')

doc.save('test-materials/database_design.docx')

# 2. Create PPTX file
prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "알고리즘과 자료구조"
subtitle.text = "컴퓨터 과학의 핵심 개념"

# Content slides
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "배열과 리스트"
content.text = "• 배열(Array): 고정 크기, O(1) 접근\n• 연결 리스트(Linked List): 동적 크기, O(n) 접근\n• 동적 배열(Dynamic Array): ArrayList, Vector"

slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "스택과 큐"
content.text = "• 스택(Stack): LIFO(Last In First Out)\n  - push(), pop(), peek()\n• 큐(Queue): FIFO(First In First Out)\n  - enqueue(), dequeue()"

slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "트리 구조"
content.text = "• 이진 트리(Binary Tree)\n• 이진 탐색 트리(BST)\n• AVL 트리\n• 레드-블랙 트리\n• B-트리"

slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "정렬 알고리즘"
content.text = "• 버블 정렬: O(n²)\n• 선택 정렬: O(n²)\n• 삽입 정렬: O(n²)\n• 퀵 정렬: O(n log n)\n• 병합 정렬: O(n log n)\n• 힙 정렬: O(n log n)"

slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "그래프 알고리즘"
content.text = "• DFS(깊이 우선 탐색)\n• BFS(너비 우선 탐색)\n• 다익스트라 알고리즘\n• 플로이드-워셜 알고리즘\n• 크루스칼 알고리즘"

prs.save('test-materials/algorithms_and_data_structures.pptx')

# 3. Create PDF file
c = canvas.Canvas("test-materials/spring_framework.pdf", pagesize=letter)
width, height = letter

# Page 1
c.setFont("Helvetica-Bold", 24)
c.drawString(100, height - 100, "Spring Framework 핵심 개념")
c.setFont("Helvetica", 12)
c.drawString(100, height - 150, "엔터프라이즈 자바 애플리케이션 개발 프레임워크")

c.setFont("Helvetica-Bold", 16)
c.drawString(100, height - 200, "1. IoC (Inversion of Control)")
c.setFont("Helvetica", 12)
c.drawString(100, height - 230, "- 객체의 생성과 의존관계 설정을 프레임워크가 담당")
c.drawString(100, height - 250, "- 개발자는 비즈니스 로직에만 집중")
c.drawString(100, height - 270, "- 느슨한 결합(Loose Coupling) 실현")

c.setFont("Helvetica-Bold", 16)
c.drawString(100, height - 320, "2. DI (Dependency Injection)")
c.setFont("Helvetica", 12)
c.drawString(100, height - 350, "- 생성자 주입 (Constructor Injection)")
c.drawString(100, height - 370, "- 세터 주입 (Setter Injection)")
c.drawString(100, height - 390, "- 필드 주입 (Field Injection)")
c.drawString(100, height - 410, "- @Autowired, @Resource, @Inject 어노테이션")

c.setFont("Helvetica-Bold", 16)
c.drawString(100, height - 460, "3. AOP (Aspect Oriented Programming)")
c.setFont("Helvetica", 12)
c.drawString(100, height - 490, "- 횡단 관심사(Cross-cutting Concerns) 분리")
c.drawString(100, height - 510, "- 로깅, 트랜잭션, 보안 등")
c.drawString(100, height - 530, "- @Aspect, @Before, @After, @Around")

# Page 2
c.showPage()
c.setFont("Helvetica-Bold", 16)
c.drawString(100, height - 100, "4. Spring MVC")
c.setFont("Helvetica", 12)
c.drawString(100, height - 130, "- Model-View-Controller 패턴")
c.drawString(100, height - 150, "- DispatcherServlet이 중앙 제어")
c.drawString(100, height - 170, "- @Controller, @RequestMapping, @ResponseBody")
c.drawString(100, height - 190, "- ViewResolver를 통한 뷰 처리")

c.setFont("Helvetica-Bold", 16)
c.drawString(100, height - 240, "5. Spring Boot")
c.setFont("Helvetica", 12)
c.drawString(100, height - 270, "- 자동 설정(Auto Configuration)")
c.drawString(100, height - 290, "- 내장 서버(Embedded Server)")
c.drawString(100, height - 310, "- 스타터 의존성(Starter Dependencies)")
c.drawString(100, height - 330, "- 운영 준비 기능(Production-ready features)")

c.setFont("Helvetica-Bold", 16)
c.drawString(100, height - 380, "6. Spring Data JPA")
c.setFont("Helvetica", 12)
c.drawString(100, height - 410, "- Repository 인터페이스")
c.drawString(100, height - 430, "- 쿼리 메소드 자동 생성")
c.drawString(100, height - 450, "- @Query 어노테이션")
c.drawString(100, height - 470, "- Pageable, Sort 지원")

c.save()

print("Test files created successfully!")
print("- database_design.docx")
print("- algorithms_and_data_structures.pptx")
print("- spring_framework.pdf")
